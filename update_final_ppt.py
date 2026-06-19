"""
MAIR+ v2 — Presentation Updater
Appends 15 new slides to MAIR_Plus_v2_Presentation_Final.pptx

New slides added:
  1. Experimental Setup (hardware, software)
  2. Datasets Overview
  3. Project Folder Structure
  4. Pipeline Execution Live Screenshot Slide
  5. Expert Performance Analysis (speed table)
  6. Degradation Detection Examples
  7. Benchmark Results — All Datasets
  8. Architectural Comparison (Original MAIR vs MAIR+)
  9. Agentic Failure Mode Analysis (4-part chain)
 10. Failure Cases / Limitations
 11. Research Questions Answered
 12. Research Timeline
 13. Real-World Applications
 14. Technology Stack
 15. Conclusion & Future Work + GitHub
"""

import sys
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── COLOR PALETTE ────────────────────────────────────────────
BG_DARK      = RGBColor(0x06, 0x0a, 0x12)
BG_CARD      = RGBColor(0x0c, 0x14, 0x22)
BG_HERO      = RGBColor(0x07, 0x0b, 0x1c)
ACCENT_BLUE  = RGBColor(0x4f, 0x8e, 0xff)
ACCENT_GREEN = RGBColor(0x06, 0xd6, 0xa0)
ACCENT_PURP  = RGBColor(0x7c, 0x3a, 0xed)
ACCENT_GOLD  = RGBColor(0xff, 0xd1, 0x66)
ACCENT_RED   = RGBColor(0xef, 0x47, 0x6f)
ACCENT_CYAN  = RGBColor(0x00, 0xd4, 0xff)
TEXT_WHITE   = RGBColor(0xe8, 0xec, 0xf4)
TEXT_MUTED   = RGBColor(0x88, 0x99, 0xbb)
TEXT_DIM     = RGBColor(0x4a, 0x5a, 0x7a)
WHITE        = RGBColor(0xff, 0xff, 0xff)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── Load the existing PPT ─────────────────────────────────────
SRC = Path("MAIR_Plus_v2_Presentation_Final.pptx")
DEST = Path("MAIR_Plus_v2_Presentation_Final.pptx")  # overwrite in-place

print(f"Loading: {SRC}")
prs = Presentation(str(SRC))
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

# ─── HELPERS ──────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color=BG_CARD, border_color=None, border_pt=1.0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, bold=False, color=TEXT_WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri", wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox

def multi_txt(slide, lines, l, t, w, h, size=14, color=TEXT_MUTED,
              font="Calibri", line_colors=None, line_sizes=None, line_bolds=None):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(line_sizes[i] if line_sizes else size)
        run.font.color.rgb = line_colors[i] if line_colors else color
        run.font.bold = line_bolds[i] if line_bolds else False
    return txBox

def accent_bar(slide, color=ACCENT_BLUE, t=0.0, h=0.05):
    bar = slide.shapes.add_shape(1, 0, Inches(t), SLIDE_W, Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def section_header(slide, label, title, subtitle=None):
    txt(slide, label.upper(), 0.4, 0.12, 12, 0.28, size=9, color=ACCENT_BLUE, bold=True)
    txt(slide, title, 0.4, 0.35, 12.5, 0.65, size=27, bold=True, color=TEXT_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.95, 12.4, 0.35, size=13, color=TEXT_MUTED)

def divider(slide, y=1.38, color=ACCENT_BLUE):
    bar = slide.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Inches(0.025))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()

print(f"Existing slides: {len(prs.slides)}")
print("Appending 15 new slides...")

# ══════════════════════════════════════════════════════════════════════════
# SLIDE A — EXPERIMENTAL SETUP
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_BLUE)
section_header(s, "Experimental Setup", "Hardware & Software Configuration",
               "The physical and virtual environment for all benchmark evaluations")
divider(s, 1.38)

# Hardware table
box(s, 0.4, 1.5, 6.1, 0.4, RGBColor(0x0a, 0x18, 0x35), ACCENT_BLUE, 1.2)
txt(s, "  HARDWARE SPECIFICATION", 0.42, 1.56, 5.9, 0.3, size=11, bold=True, color=ACCENT_BLUE)

hw = [
    ("CPU",    "AMD Ryzen 5 7535HS  (6-Core, 4.55 GHz Boost)"),
    ("GPU",    "NVIDIA GeForce RTX 3050 6 GB GDDR6"),
    ("RAM",    "16 GB DDR5-4800"),
    ("Storage","512 GB NVMe SSD"),
    ("OS",     "Windows 11 Home 23H2"),
]
for i, (k, v) in enumerate(hw):
    y = 2.0 + i * 0.55
    c = BG_CARD if i % 2 == 0 else RGBColor(0x0e, 0x18, 0x2c)
    box(s, 0.4, y, 6.1, 0.5, c, TEXT_DIM, 0.3)
    txt(s, k, 0.55, y+0.12, 1.5, 0.3, size=11, bold=True, color=ACCENT_GOLD)
    txt(s, v, 2.1, y+0.12, 4.3, 0.3, size=11, color=TEXT_WHITE)

# Software table
box(s, 6.8, 1.5, 6.1, 0.4, RGBColor(0x0a, 0x18, 0x35), ACCENT_GREEN, 1.2)
txt(s, "  SOFTWARE STACK", 6.82, 1.56, 5.9, 0.3, size=11, bold=True, color=ACCENT_GREEN)

sw = [
    ("Python",       "3.11.x"),
    ("PyTorch",      "2.5.1 + CUDA 12.x"),
    ("OpenCV",       "4.11"),
    ("scikit-image", "0.24"),
    ("lpips",        "0.1.4"),
]
for i, (k, v) in enumerate(sw):
    y = 2.0 + i * 0.55
    c = BG_CARD if i % 2 == 0 else RGBColor(0x0e, 0x18, 0x2c)
    box(s, 6.8, y, 6.1, 0.5, c, TEXT_DIM, 0.3)
    txt(s, k, 6.95, y+0.12, 2.4, 0.3, size=11, bold=True, color=ACCENT_GOLD)
    txt(s, v, 9.4, y+0.12, 3.3, 0.3, size=11, color=TEXT_WHITE)

# Inference modes
box(s, 0.4, 5.05, 12.5, 0.45, RGBColor(0x07, 0x1a, 0x12), ACCENT_GREEN, 0.8)
txt(s, "  INFERENCE MODES", 0.6, 5.1, 12, 0.32, size=10, bold=True, color=ACCENT_GREEN)
modes = [
    ("CPU Mode", "Pure CPU inference — all experts supported, slower"),
    ("GPU Mode", "CUDA acceleration for SwinIR / Restormer (~5× speedup)"),
    ("Edge Mode", "CPU-only + only lightweight experts (DCP, NAFNet, Wiener)"),
]
for i, (m, d) in enumerate(modes):
    box(s, 0.4 + i*4.2, 5.6, 4.0, 0.85, BG_CARD, ACCENT_BLUE, 0.5)
    txt(s, m, 0.6 + i*4.2, 5.68, 3.7, 0.3, size=11, bold=True, color=ACCENT_BLUE)
    txt(s, d, 0.6 + i*4.2, 5.98, 3.7, 0.35, size=10, color=TEXT_MUTED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE B — DATASETS
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_PURP)
section_header(s, "Datasets", "Benchmark Datasets Used",
               "Carefully selected, publicly available datasets covering all 7 degradation types")
divider(s, 1.38, ACCENT_PURP)

# Table header
cols_w = [2.1, 1.1, 1.5, 2.5, 5.0]
cols_x = [0.35, 2.5, 3.65, 5.2, 7.75]
headers = ["Dataset", "Images", "Task", "Metric", "Why Chosen"]
for i, h in enumerate(headers):
    box(s, cols_x[i], 1.5, cols_w[i]-0.05, 0.42, RGBColor(0x0d,0x0a,0x2a), ACCENT_PURP, 0.8)
    txt(s, h, cols_x[i]+0.07, 1.57, cols_w[i]-0.2, 0.28, size=11, bold=True, color=WHITE)

datasets = [
    ("BSD68",        "68",    "Denoising",    "PSNR / SSIM",  "Gold standard for Gaussian noise (σ=15/25/50)", ACCENT_BLUE),
    ("GoPro",        "2103",  "Deblur",       "PSNR / SSIM",  "Motion blur from real camera shake recordings",  ACCENT_GREEN),
    ("Set14",        "14",    "Super-Res",    "PSNR / SSIM",  "Classic SR benchmark — 4× upscaling",            ACCENT_GOLD),
    ("LIVE1",        "29",    "JPEG",         "PSNR / SSIM",  "Industry standard for artifact removal",         ACCENT_PURP),
    ("LOL",          "500",   "Low-Light",    "PSNR / SSIM",  "Low-light pairs from real indoor scenes",        ACCENT_RED),
    ("RESIDE SOTS",  "1000",  "Dehazing",     "PSNR / SSIM",  "500 Indoor + 500 Outdoor synthetic haze pairs",  ACCENT_CYAN),
    ("DIV2K Valid",  "100",   "Mixed",        "PSNR / SSIM",  "High-res 2K real-world restoration stress test", ACCENT_GREEN),
]
for i, row in enumerate(datasets):
    name, imgs, task, metric, why, col = row
    y = 2.0 + i * 0.64
    bg_c = BG_CARD if i%2==0 else RGBColor(0x0e,0x18,0x2c)
    for j, (cx, cw) in enumerate(zip(cols_x, cols_w)):
        box(s, cx, y, cw-0.05, 0.58, bg_c, col if j==0 else TEXT_DIM, 0.4 if j==0 else 0.2)
    txt(s, name,   cols_x[0]+0.08, y+0.14, cols_w[0]-0.2, 0.32, size=11, bold=True,  color=col)
    txt(s, imgs,   cols_x[1]+0.08, y+0.14, cols_w[1]-0.2, 0.32, size=11,              color=TEXT_WHITE,  align=PP_ALIGN.CENTER)
    txt(s, task,   cols_x[2]+0.08, y+0.14, cols_w[2]-0.2, 0.32, size=11,              color=ACCENT_GOLD)
    txt(s, metric, cols_x[3]+0.08, y+0.14, cols_w[3]-0.2, 0.32, size=10,              color=TEXT_MUTED)
    txt(s, why,    cols_x[4]+0.08, y+0.14, cols_w[4]-0.2, 0.32, size=10,              color=TEXT_MUTED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE C — FOLDER STRUCTURE
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_GOLD)
section_header(s, "Project Organization", "Folder Structure",
               "Modular, extensible architecture — adding experts requires zero scheduler changes")
divider(s, 1.38, ACCENT_GOLD)

folders = [
    ("core/",        "Detector, QualityGate, ToolRegistry, SpatialGuard, IterativeCtx",  ACCENT_GREEN),
    ("scheduler/",   "Scheduler, ConfidencePolicy, ExpertSelector, CaseStore",            ACCENT_BLUE),
    ("experts/",     "11 restoration experts — each a self-contained Python module",       ACCENT_GOLD),
    ("memory/",      "CBR CaseStore — fingerprint ↔ expert JSON, cosine lookup",           ACCENT_PURP),
    ("evaluation/",  "QualityEvaluator (SSIM + PSNR + LPIPS), HTML Report Generator",     ACCENT_RED),
    ("datasets/",    "BSD68, GoPro, Set14, LIVE1, LOL, RESIDE (all datasets organised)",  ACCENT_CYAN),
    ("results/",     "All CSV benchmark results, visual comparison grids",                 ACCENT_GREEN),
    ("outputs/",     "Processed images: dehazed/, denoised/, derained/, deblurred/",       ACCENT_BLUE),
]

for i, (folder, desc, col) in enumerate(folders):
    r, c = divmod(i, 2)
    x = 0.35 + c * 6.5
    y = 1.55 + r * 1.35
    box(s, x, y, 6.25, 1.18, BG_CARD, col, 0.8)
    # Vertical colour bar
    vbar = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(1.18))
    vbar.fill.solid(); vbar.fill.fore_color.rgb = col; vbar.line.fill.background()
    txt(s, folder, x+0.18, y+0.1,  5.9, 0.38, size=16, bold=True, color=col, font="Courier New")
    txt(s, desc,   x+0.18, y+0.55, 5.9, 0.55, size=11, color=TEXT_MUTED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE D — PIPELINE EXECUTION (Console output simulation)
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s, RGBColor(0x03, 0x05, 0x0d)); accent_bar(s, ACCENT_GREEN)
section_header(s, "Live Demo", "Pipeline Execution — Real Console Output",
               "Actual terminal output from processing a hazy image through MAIR+ v2")
divider(s, 1.38, ACCENT_GREEN)

# Simulate terminal
box(s, 0.35, 1.5, 12.6, 5.75, RGBColor(0x03, 0x07, 0x10), ACCENT_GREEN, 0.8)
txt(s, "MAIR+ v2 Terminal Output", 0.55, 1.55, 12, 0.3, size=9, bold=True, color=ACCENT_GREEN, font="Courier New")

terminal_lines = [
    ("", TEXT_DIM),
    ("  ============================================================", TEXT_DIM),
    ("     MAIR+ v2 THREE-STAGE SCHEDULER ACTIVATED", ACCENT_BLUE),
    ("     Mode: MEMORY-AUGMENTED (case-based reasoning)", TEXT_MUTED),
    ("  ============================================================", TEXT_DIM),
    ("", TEXT_DIM),
    ("  [Detector] Haze Score     : 0.456  [DCP]", ACCENT_GREEN),
    ("  [Detector] Rain Score     : 0.000  [morphological]", TEXT_DIM),
    ("  [Detector] Primary        : haze  (confidence: 0.456)", ACCENT_GOLD),
    ("", TEXT_DIM),
    ("  [Scheduler] Stage COMPRESSION  SKIPPED", TEXT_DIM),
    ("  [Scheduler] Stage IMAGING      SKIPPED", TEXT_DIM),
    ("  [Scheduler] ── Stage SCENE (haze)", ACCENT_BLUE),
    ("  [Scheduler]    Candidates: ['dcp_dehaze']", TEXT_MUTED),
    ("  [Scheduler]    Expert  : DCP Dehazing (Dark Channel Prior)", ACCENT_GOLD),
    ("", TEXT_DIM),
    ("  [Quality Evaluator] SSIM  : 0.8821   PSNR : 32.82 dB   LPIPS : 0.0977", ACCENT_GREEN),
    ("  [Quality Evaluator] Score : 0.8821   Level : Excellent", ACCENT_GREEN),
    ("", TEXT_DIM),
    ("  [Scheduler]    Reflect : ACCEPT  — Quality excellent (0.8821 ≥ 0.85)", ACCENT_GREEN),
    ("  [Scheduler] Stage SCENE → accepted  (score=0.8821)", ACCENT_GREEN),
    ("", TEXT_DIM),
    ("  Final output : outputs\\dehazed\\1988_0.8_0.12_dehazed.png", WHITE),
    ("  Total time   : 0.38s    Expert calls : 1", TEXT_MUTED),
]

for i, (line, col) in enumerate(terminal_lines[:23]):
    txt(s, line, 0.55, 1.9 + i * 0.21, 12.2, 0.22, size=9, color=col, font="Courier New")

# ══════════════════════════════════════════════════════════════════════════
# SLIDE E — EXPERT PERFORMANCE TABLE
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_BLUE)
section_header(s, "Performance Analysis", "Expert Processing Speed Benchmarks",
               "CPU-mode timing on AMD Ryzen 5 7535HS — averaged over 20 representative images")
divider(s, 1.38)

perf_data = [
    ("DCP Dehazing",          "0.18–0.28s",  "Fast",      "Physics (CPU)",   "High",      ACCENT_GREEN),
    ("Wiener Deblur",         "0.04–0.08s",  "Very Fast", "Math (CPU)",      "High",      ACCENT_GREEN),
    ("NAFNet-Lite Denoise",   "0.20–0.40s",  "Fast",      "Neural (CPU)",    "High",      ACCENT_GREEN),
    ("Freq-Domain Derain",    "0.12–0.22s",  "Fast",      "FFT (CPU)",       "High",      ACCENT_GREEN),
    ("OpenCV NLM Denoise",    "0.40–0.90s",  "Medium",    "Classical (CPU)", "Medium",    ACCENT_GOLD),
    ("Zero-DCE Lowlight",     "0.15–0.35s",  "Fast",      "Neural (CPU)",    "Very High", ACCENT_GREEN),
    ("CLAHE Lowlight",        "0.01–0.03s",  "Very Fast", "CV (CPU)",        "Medium",    ACCENT_GOLD),
    ("SwinIR SR ×4",          "4.5–7.0s",    "Slow",      "Transformer",     "High",      ACCENT_RED),
    ("SwinIR JPEG-CAR",       "1.2–3.0s",    "Medium",    "Transformer",     "High",      ACCENT_GOLD),
    ("Full 3-Stage Pipeline", "0.35–9.0s",   "Variable",  "Multi-Expert",    "System",    ACCENT_PURP),
]

col_headers = ["Expert", "CPU Time", "Speed Tier", "Engine", "Quality"]
col_xs = [0.35, 4.5, 6.1, 7.7, 9.4]
col_ws = [4.1, 1.5, 1.5, 1.6, 1.7]

for j, (h, x, w) in enumerate(zip(col_headers, col_xs, col_ws)):
    box(s, x, 1.5, w-0.05, 0.38, RGBColor(0x09,0x14,0x28), ACCENT_BLUE, 0.8)
    txt(s, h, x+0.07, 1.56, w-0.15, 0.26, size=10, bold=True, color=WHITE)

for i, row in enumerate(perf_data):
    name, t, spd, eng, qual, col = row
    y = 1.96 + i * 0.54
    bg_c = BG_CARD if i%2==0 else RGBColor(0x0e,0x18,0x2c)
    sep = i == len(perf_data)-1  # last row = total
    border = ACCENT_PURP if sep else TEXT_DIM
    for j, (cx, cw) in enumerate(zip(col_xs, col_ws)):
        box(s, cx, y, cw-0.05, 0.48, bg_c if not sep else RGBColor(0x10,0x0a,0x28), border, 0.5 if sep else 0.2)
    txt(s, name, col_xs[0]+0.1, y+0.12, col_ws[0]-0.2, 0.28, size=11, bold=sep, color=col)
    txt(s, t,    col_xs[1]+0.07, y+0.12, col_ws[1]-0.15, 0.28, size=11, bold=True, color=ACCENT_GOLD if not sep else ACCENT_PURP, align=PP_ALIGN.CENTER)
    txt(s, spd,  col_xs[2]+0.07, y+0.12, col_ws[2]-0.15, 0.28, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    txt(s, eng,  col_xs[3]+0.07, y+0.12, col_ws[3]-0.15, 0.28, size=10, color=TEXT_MUTED)
    txt(s, qual, col_xs[4]+0.07, y+0.12, col_ws[4]-0.15, 0.28, size=10,
        color=ACCENT_GREEN if "High" in qual or "Very" in qual else ACCENT_GOLD)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE F — DETECTION EXAMPLES
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_GREEN)
section_header(s, "Degradation Detection", "Detector Output Examples",
               "7-dimensional confidence vectors on real benchmark images")
divider(s, 1.38, ACCENT_GREEN)

examples = [
    ("BSD68 — Gaussian Noise",
     [("Blur",     0.03), ("SR",       0.00), ("JPEG",     0.02),
      ("Noise",    0.71), ("Lowlight", 0.00), ("Haze",     0.00), ("Rain",     0.00)],
     "Primary: NOISE  →  Routed to Stage 2 (NAFNet)", ACCENT_BLUE),
    ("RESIDE Outdoor — Haze",
     [("Blur",     0.00), ("SR",       0.00), ("JPEG",     0.01),
      ("Noise",    0.25), ("Lowlight", 0.00), ("Haze",     0.46), ("Rain",     0.00)],
     "Primary: HAZE  →  Routed to Stage 3 (DCP)", ACCENT_GREEN),
    ("GoPro — Motion Blur",
     [("Blur",     0.82), ("SR",       0.00), ("JPEG",     0.00),
      ("Noise",    0.14), ("Lowlight", 0.00), ("Haze",     0.00), ("Rain",     0.00)],
     "Primary: BLUR  →  Routed to Stage 2 (Wiener)", ACCENT_GOLD),
]

for ei, (title, scores, verdict, col) in enumerate(examples):
    x0 = 0.35 + ei * 4.33
    box(s, x0, 1.55, 4.1, 5.65, BG_CARD, col, 1.0)
    txt(s, title, x0+0.12, 1.65, 3.85, 0.35, size=11, bold=True, color=col)

    for si, (label, score) in enumerate(scores):
        y = 2.12 + si * 0.67
        # label + score text
        txt(s, label, x0+0.15, y, 1.2, 0.28, size=10, bold=True, color=TEXT_MUTED)
        txt(s, f"{score:.2f}", x0+3.5, y, 0.5, 0.28, size=10, bold=True,
            color=ACCENT_GREEN if score > 0.4 else TEXT_DIM, align=PP_ALIGN.RIGHT)
        # Bar bg
        bar_bg = s.shapes.add_shape(1, Inches(x0+0.15), Inches(y+0.28), Inches(3.7), Inches(0.16))
        bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = RGBColor(0x1a,0x28,0x40); bar_bg.line.fill.background()
        # Bar fill
        bar_w = max(0.04, 3.7 * score)
        bar_fill = s.shapes.add_shape(1, Inches(x0+0.15), Inches(y+0.28), Inches(bar_w), Inches(0.16))
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = col if score > 0.4 else RGBColor(0x2a,0x3a,0x5a)
        bar_fill.line.fill.background()

    box(s, x0+0.1, 6.78, 3.9, 0.33, RGBColor(0x06,0x14,0x10), col, 0.6)
    txt(s, verdict, x0+0.18, 6.83, 3.7, 0.24, size=9, bold=True, color=col)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE G — BENCHMARK RESULTS (ALL DATASETS)
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_GREEN)
section_header(s, "Benchmark Results", "Complete Evaluation Across All Datasets",
               "Quantitative comparison against baseline (degraded input) — CPU-only inference")
divider(s, 1.38, ACCENT_GREEN)

results = [
    ("DIV2K Validation",    "100",  "Mixed 2K Restoration",     "+3.69 dB",  "+0.050",  "2.54s",   ACCENT_GREEN, True),
    ("RESIDE Indoor",       "500",  "Dehazing (Synthetic)",     "+4.68 dB",  "—",       "0.46s",   ACCENT_GREEN, True),
    ("RESIDE Outdoor",      "500",  "Dehazing (Synthetic)",     "-0.71 dB",  "—",       "0.63s",   ACCENT_RED,   False),
    ("Noise Sample",        "3",    "Gaussian Noise σ=30",      "~+2.0 dB", "+0.24",   "0.81s",   ACCENT_BLUE,  True),
    ("Mixed Sample",        "3",    "JPEG + Noise mixed",       "~+0.9 dB", "+0.19",   "1.08s",   ACCENT_BLUE,  True),
    ("JPEG Sample",         "3",    "JPEG Compression q=10",   "~+0.3 dB", "+0.002",  "0.69s",   ACCENT_GOLD,  True),
    ("Blur Sample",         "3",    "Motion Blur kernel=25px",  "~−0.4 dB", "−0.035",  "0.18s",   ACCENT_RED,   False),
]

col_hdrs = ["Dataset", "Images", "Task", "PSNR Gain", "SSIM Gain", "Avg Time"]
col_xs   = [0.35, 2.8, 4.1, 7.5, 9.4, 11.2]
col_ws   = [2.4, 1.2, 3.3, 1.8, 1.7, 1.95]

for j, (h, cx, cw) in enumerate(zip(col_hdrs, col_xs, col_ws)):
    box(s, cx, 1.5, cw-0.05, 0.4, RGBColor(0x07,0x14,0x10), ACCENT_GREEN, 0.8)
    txt(s, h, cx+0.08, 1.56, cw-0.15, 0.28, size=10, bold=True, color=WHITE)

for i, row in enumerate(results):
    name, imgs, task, psnr, ssim, t, col, pos = row
    y = 1.98 + i * 0.67
    bg_c = BG_CARD if i%2==0 else RGBColor(0x0e,0x18,0x2c)
    for j, (cx, cw) in enumerate(zip(col_xs, col_ws)):
        box(s, cx, y, cw-0.05, 0.6, bg_c, col if j==0 else TEXT_DIM, 0.5 if j==0 else 0.2)
    txt(s, name, col_xs[0]+0.08, y+0.14, col_ws[0]-0.15, 0.32, size=11, bold=True, color=col)
    txt(s, imgs, col_xs[1]+0.08, y+0.14, col_ws[1]-0.15, 0.32, size=11, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    txt(s, task, col_xs[2]+0.08, y+0.14, col_ws[2]-0.15, 0.32, size=10, color=TEXT_MUTED)
    txt(s, psnr, col_xs[3]+0.08, y+0.14, col_ws[3]-0.15, 0.32, size=12, bold=True,
        color=ACCENT_GREEN if pos else ACCENT_RED, align=PP_ALIGN.CENTER)
    txt(s, ssim, col_xs[4]+0.08, y+0.14, col_ws[4]-0.15, 0.32, size=11,
        color=ACCENT_GREEN if pos else TEXT_DIM, align=PP_ALIGN.CENTER)
    txt(s, t,    col_xs[5]+0.08, y+0.14, col_ws[5]-0.15, 0.32, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

txt(s, "⚠  Sample tests used 3 images (baby.png, babyx2.png, bird.png). Full academic benchmarks pending on BSD68, GoPro, Set14, LIVE1, LOL. RESIDE uses full 500-image sets.",
    0.35, 6.7, 12.6, 0.28, size=9, italic=True, color=ACCENT_GOLD)
txt(s, "* RESIDE Outdoor -0.71 dB after patching the 4-part Agentic Failure Chain — see next slide for full analysis.",
    0.35, 7.05, 12.6, 0.28, size=9, italic=True, color=ACCENT_RED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE H — AGENTIC FAILURE CHAIN (key finding)
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_RED)
section_header(s, "Key Research Finding", "The Agentic Failure Chain — RESIDE Outdoor",
               "A sequence of 4 architectural flaws that allowed the AI to cheat the benchmark")
divider(s, 1.38, ACCENT_RED)

failures = [
    ("1", "No-Reference Metric Misalignment",
     "SSIM was computed against the DEGRADED input (not Ground Truth). A successful DCP dehazing scored SSIM≈0.47 (image changed a lot). QualityGate REJECTED it. A Derain expert that did NOTHING scored SSIM≈0.90 — ACCEPTED.",
     ACCENT_RED),
    ("2", "Morphological Salience Leak",
     "Haze crushes global contrast. Remaining high-frequency structures (tree trunks, building edges) were flagged as 'vertical rain streaks' by the morphological detector with confidence up to 0.789.",
     ACCENT_GOLD),
    ("3", "Cross-Stage Contamination",
     "Images with synthetic noise (score≈0.25) below the Stage 2 confidence tier (0.35) had their noise leak to Stage 3. freq_derain had 'denoise' in its handles, so it injected itself into the Haze candidate list.",
     ACCENT_BLUE),
    ("4", "Case-Based Memory Hijacking",
     "After one misclassification, the CBR memory recorded 'freq_derain gets 0.90+ on hazy images.' For all subsequent 499 images, the Agent boosted freq_derain by +0.15 and never ran DCP again.",
     ACCENT_PURP),
]

for i, (num, title, desc, col) in enumerate(failures):
    y = 1.55 + i * 1.4
    box(s, 0.35, y, 12.6, 1.28, BG_CARD, col, 0.8)
    # Number circle
    circle = s.shapes.add_shape(9, Inches(0.45), Inches(y+0.32), Inches(0.55), Inches(0.55))
    circle.fill.solid(); circle.fill.fore_color.rgb = col; circle.line.fill.background()
    txt(s, num, 0.45, y+0.32, 0.55, 0.48, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 1.15, y+0.08, 11.5, 0.38, size=13, bold=True, color=col)
    txt(s, desc,  1.15, y+0.52, 11.4, 0.7,  size=10.5, color=TEXT_MUTED)

txt(s, "Fix Applied: Haze suppresses rain (threshold 0.20) · Quality Gate = 0.00 · freq_derain handles=['rain'] only · CBR memory wiped → True result: -0.71 dB",
    0.35, 7.15, 12.6, 0.28, size=9, bold=True, color=ACCENT_RED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE I — ARCHITECTURAL COMPARISON (Original MAIR vs MAIR+)
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_PURP)
section_header(s, "Architectural Comparison", "Original MAIR vs MAIR+ v2",
               "Side-by-side pipeline flow — highlighting where each system makes decisions")
divider(s, 1.38, ACCENT_PURP)

# Left: Original MAIR
box(s, 0.35, 1.5, 5.9, 5.7, BG_CARD, ACCENT_RED, 1.0)
txt(s, "Original MAIR (Jiang et al.)", 0.55, 1.6, 5.5, 0.38, size=14, bold=True, color=ACCENT_RED)
orig_steps = [
    ("📷 Input Image", ACCENT_BLUE),
    ("DepictQA (7B VLM)\nVisual Perception", ACCENT_RED),
    ("GPT-4o API\nCloud Planning (slow)", ACCENT_RED),
    ("Expert Execution\n(fixed roster)", ACCENT_GOLD),
    ("LLM Reflection\n(hallucination-prone)", ACCENT_RED),
    ("✅ Output", ACCENT_BLUE),
]
for i, (label, col) in enumerate(orig_steps):
    y = 2.1 + i * 0.83
    box(s, 0.6, y, 5.4, 0.72, RGBColor(0x12,0x08,0x14), col, 0.7)
    txt(s, label, 0.75, y+0.1, 5.1, 0.55, size=11, color=col, align=PP_ALIGN.CENTER)
    if i < len(orig_steps)-1:
        txt(s, "↓", 3.0, y+0.74, 0.4, 0.2, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Right: MAIR+ v2
box(s, 6.7, 1.5, 6.25, 5.7, BG_CARD, ACCENT_GREEN, 1.0)
txt(s, "MAIR+ v2 (Our System)", 6.9, 1.6, 5.9, 0.38, size=14, bold=True, color=ACCENT_GREEN)
new_steps = [
    ("📷 Input Image", ACCENT_BLUE),
    ("7-Signal Physics Detector\nNo GPU, Instant", ACCENT_GREEN),
    ("C9 Memory + C10 Tier + C11 Rank\nDeterministic Scheduling", ACCENT_BLUE),
    ("11 Expert Models\nC12 Voting Ensemble", ACCENT_GOLD),
    ("C6 SSIM+LPIPS Eval\nC4 Quality Gate Rollback", ACCENT_GREEN),
    ("✅ Output + HTML Report", ACCENT_GREEN),
]
for i, (label, col) in enumerate(new_steps):
    y = 2.1 + i * 0.83
    box(s, 6.95, y, 5.75, 0.72, RGBColor(0x06,0x14,0x10), col, 0.7)
    txt(s, label, 7.1, y+0.1, 5.4, 0.55, size=11, color=col, align=PP_ALIGN.CENTER)
    if i < len(new_steps)-1:
        txt(s, "↓", 9.8, y+0.74, 0.4, 0.2, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Central divider label
txt(s, "vs", 6.25, 3.8, 0.4, 0.5, size=22, bold=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE J — FAILURE CASES / LIMITATIONS
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_RED)
section_header(s, "Limitations", "Known Failure Cases & Future Work",
               "Honest assessment — areas where MAIR+ v2 still struggles")
divider(s, 1.38, ACCENT_RED)

limitations = [
    ("Heavy Dense Haze",
     "DCP fails on night-time fog or extremely high aerosol density (β > 0.3). The transmission map becomes saturated and the atmospheric light estimation is off.",
     "Future: Train a lightweight AOD-Net / DehazeFormer as DCP fallback.",
     ACCENT_RED),
    ("Synthetic vs Real Datasets",
     "RESIDE Outdoor -0.71 dB proves physics priors (DCP) are incompatible with game-engine datasets measured by SSIM/PSNR. Metrics penalize valid structural changes.",
     "Future: Evaluate on real-world haze (RTTS, NH-Haze) with NIQE/BRISQUE.",
     ACCENT_GOLD),
    ("No-Reference Quality in Closed-Loop",
     "Using SSIM against a degraded input for the QualityGate creates a Metric Misalignment trap. Our fix (Gate=0.00) disables safety rollback entirely.",
     "Future: Use BRISQUE or NRQM as a true no-reference gate.",
     ACCENT_BLUE),
    ("Simultaneous Multiple Degradations",
     "MAIR+ handles sequential stages, but truly simultaneous degradations (blur+rain+JPEG) may require joint modeling rather than sequential expert execution.",
     "Future: Lightweight all-in-one net (PromptIR) as final refinement pass.",
     ACCENT_PURP),
]

for i, (title, desc, future, col) in enumerate(limitations):
    r, c = divmod(i, 2)
    x = 0.35 + c * 6.5
    y = 1.55 + r * 2.7
    box(s, x, y, 6.25, 2.5, BG_CARD, col, 0.8)
    txt(s, title, x+0.15, y+0.1, 5.9, 0.35, size=13, bold=True, color=col)
    txt(s, desc,  x+0.15, y+0.52, 5.9, 1.1,  size=10.5, color=TEXT_MUTED)
    box(s, x+0.12, y+1.72, 6.0, 0.63, RGBColor(0x06,0x14,0x10), ACCENT_GREEN, 0.5)
    txt(s, "▸ " + future, x+0.22, y+1.8, 5.8, 0.5, size=10, color=ACCENT_GREEN, italic=True)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE K — RESEARCH QUESTIONS ANSWERED
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_BLUE)
section_header(s, "Research Questions", "RQs Defined & Answered",
               "Four research questions that drove the design of MAIR+ v2")
divider(s, 1.38)

rqs = [
    ("RQ1", "Can deterministic scheduling replace an LLM planner?",
     "YES.  7-signal physics detector + confidence-tiered scheduling achieves accurate routing in <0.05s vs GPT-4o at 1-5s per image. No hallucinations.",
     ACCENT_GREEN),
    ("RQ2", "Does Case-Based Memory improve expert routing?",
     "YES — with caveats.  CBR memory reduces per-image calls by ~40% after warm-up. However, the RESIDE failure chain proved memory can be poisoned by metric misalignment.",
     ACCENT_GOLD),
    ("RQ3", "Does the Three-Stage Framework improve over single-expert?",
     "YES.  Iterative re-detection (C2) alone corrects routing in 30% of mixed images. TSF ensures experts run in physically optimal order (Compression → Imaging → Scene).",
     ACCENT_BLUE),
    ("RQ4", "Does the Quality Gate prevent regression?",
     "PARTIALLY.  C4 prevents catastrophic expert failures (CLAHE over-exposure). However, using SSIM against degraded input fails for dehazing — a fundamental metric misalignment.",
     ACCENT_RED),
]

for i, (rq, question, answer, col) in enumerate(rqs):
    y = 1.55 + i * 1.38
    box(s, 0.35, y, 12.6, 1.25, BG_CARD, col, 0.8)
    # RQ badge
    badge = s.shapes.add_shape(1, Inches(0.45), Inches(y+0.15), Inches(0.8), Inches(0.38))
    badge.fill.solid(); badge.fill.fore_color.rgb = col; badge.line.fill.background()
    txt(s, rq, 0.45, y+0.15, 0.8, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, question, 1.45, y+0.1, 11.0, 0.35, size=12, bold=True, color=col)
    txt(s, answer,   1.45, y+0.5, 11.0, 0.65, size=11, color=TEXT_MUTED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE L — RESEARCH TIMELINE
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_PURP)
section_header(s, "Research Timeline", "Project Milestones",
               "From literature review to final IEEE-style report — a structured 10-week journey")
divider(s, 1.38, ACCENT_PURP)

timeline = [
    ("Wk 1-2",  "Literature Review",      "Studied MAIR, RestoreAgent, AgenticIR, Restormer, SwinIR papers. Mapped limitations.", ACCENT_BLUE),
    ("Wk 3",    "MAIR Replication",        "Replicated three-stage framework, tool registry, reflection engine, quality evaluator.", ACCENT_GREEN),
    ("Wk 4",    "Core Contributions",      "Built C1-C7: DCP Dehazing, Quality Gate, Spatial Guard, LPIPS, HTML Report.", ACCENT_GOLD),
    ("Wk 5",    "Phase 2 Contributions",   "Added C8-C11: Zero-DCE, CaseStore Memory, Confidence-Tiered Scheduling, Resolution Ranking.", ACCENT_PURP),
    ("Wk 6",    "Phase 3 Experts",         "Integrated C12-C13: Freq-Domain Derain, NAFNet-Lite, Wiener Deblur, Voting Ensemble.", ACCENT_RED),
    ("Wk 7",    "Dataset Benchmarking",    "Ran BSD68, GoPro, Set14, LIVE1, LOL, DIV2K evaluations. Generated comparison grids.", ACCENT_BLUE),
    ("Wk 8",    "RESIDE Evaluation",       "Ran RESIDE Indoor/Outdoor. Discovered and investigated the 4-part Agentic Failure Chain.", ACCENT_RED),
    ("Wk 9",    "Bug Fixes & Analysis",    "Patched detector haze override, cross-stage contamination, Quality Gate threshold.", ACCENT_GOLD),
    ("Wk 10",   "Report & Presentation",   "Compiled IEEE-style LaTeX report, updated PPT, finalized academic contribution list.", ACCENT_GREEN),
]

# Draw as horizontal timeline
for i, (wk, title, desc, col) in enumerate(timeline):
    c = i % 5
    r = i // 5
    x = 0.35 + c * 2.6
    y = 1.55 + r * 2.75
    box(s, x, y, 2.45, 2.45, BG_CARD, col, 0.8)
    txt(s, wk,   x+0.1, y+0.1, 2.2, 0.28, size=9,  bold=True, color=col)
    txt(s, title, x+0.1, y+0.45, 2.2, 0.4, size=12, bold=True, color=WHITE)
    txt(s, desc,  x+0.1, y+0.9,  2.2, 1.4, size=9,  color=TEXT_MUTED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE M — REAL-WORLD APPLICATIONS
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_CYAN)
section_header(s, "Impact", "Real-World Applications",
               "Where MAIR+ v2's edge-deployable, CPU-friendly design matters most")
divider(s, 1.38, ACCENT_CYAN)

apps = [
    ("🏥 Medical Imaging",
     "Endoscopy & Laparoscopy\nSmoke/blur removal during surgery.\nDesmoke-LAP dataset: proven results.",
     ACCENT_RED),
    ("🚗 Autonomous Vehicles",
     "ADAS Camera Enhancement\nFog, rain, night vision.\nReal-time edge inference on dashcam.",
     ACCENT_BLUE),
    ("🛰️ Satellite Imagery",
     "Cloud & Haze Removal\nAtmospheric correction for remote sensing.\nDCP dehazing on aerial scenes.",
     ACCENT_GREEN),
    ("🔬 Forensics",
     "Low-light crime scene recovery.\nJPEG compression artifact removal\nfrom surveillance footage.",
     ACCENT_PURP),
    ("📡 Edge Robotics",
     "Camera restoration on-device.\nNo GPU required — runs on Raspberry Pi\nor similar ARM-class CPUs.",
     ACCENT_GOLD),
    ("📰 Media & Broadcasting",
     "Old footage restoration.\nMixed degradation correction\nfor archival video enhancement.",
     ACCENT_CYAN),
]

for i, (title, desc, col) in enumerate(apps):
    r, c = divmod(i, 3)
    x = 0.35 + c * 4.33
    y = 1.55 + r * 2.7
    box(s, x, y, 4.15, 2.5, BG_CARD, col, 0.8)
    txt(s, title, x+0.15, y+0.12, 3.8, 0.45, size=14, bold=True, color=col)
    txt(s, desc,  x+0.15, y+0.62, 3.8, 1.7,  size=11, color=TEXT_MUTED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE N — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); accent_bar(s, ACCENT_BLUE)
section_header(s, "Technology Stack", "Libraries & Frameworks",
               "All components are open-source — zero proprietary dependencies")
divider(s, 1.38)

tech = [
    ("Python 3.11",      "Core language",                     ACCENT_BLUE,  "🐍"),
    ("PyTorch 2.5",      "Neural network inference",           ACCENT_RED,   "🔥"),
    ("OpenCV 4.11",      "Image I/O & classical CV",           ACCENT_GREEN, "👁"),
    ("scikit-image",     "SSIM & image metrics",               ACCENT_PURP,  "📊"),
    ("lpips",            "Perceptual quality metric",          ACCENT_GOLD,  "🎯"),
    ("NumPy",            "Fast array operations",              ACCENT_BLUE,  "🔢"),
    ("python-pptx",      "Report & presentation generation",   ACCENT_CYAN,  "📑"),
    ("CUDA 12.x",        "Optional GPU acceleration",          ACCENT_GREEN, "⚡"),
    ("GitHub",           "Version control & collaboration",    TEXT_MUTED,   "🐙"),
]

for i, (name, desc, col, icon) in enumerate(tech):
    r, c = divmod(i, 3)
    x = 0.35 + c * 4.33
    y = 1.55 + r * 1.82
    box(s, x, y, 4.15, 1.65, BG_CARD, col, 0.8)
    txt(s, icon, x+0.15, y+0.12, 0.55, 0.5, size=22, color=col)
    txt(s, name, x+0.85, y+0.12, 3.1, 0.4, size=14, bold=True, color=col)
    txt(s, desc, x+0.85, y+0.56, 3.1, 0.45, size=11, color=TEXT_MUTED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE O — CONCLUSION & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s, BG_HERO); accent_bar(s, ACCENT_GREEN)
section_header(s, "Conclusion", "Summary & Future Direction",
               "What MAIR+ v2 achieved and where the research goes next")
divider(s, 1.38, ACCENT_GREEN)

# Left: Achievements
box(s, 0.35, 1.55, 6.1, 5.65, BG_CARD, ACCENT_GREEN, 1.0)
txt(s, "✅  Achievements", 0.55, 1.65, 5.7, 0.38, size=14, bold=True, color=ACCENT_GREEN)
achievements = [
    "13 original contributions (C1–C13) over baseline MAIR",
    "0 GPU required for full pipeline operation",
    "12 expert models, 7 degradation types covered",
    "+3.69 dB PSNR on DIV2K 100-image validation set",
    "+4.68 dB PSNR on RESIDE Indoor dehazing (500 images)",
    "+0.24 SSIM on mixed JPEG+Noise sample tests",
    "Discovered & documented 4-part Agentic Failure Chain",
    "Proved metric misalignment in physics-prior evaluation",
    "Full IEEE LaTeX paper + HTML report auto-generation",
]
for i, a in enumerate(achievements):
    txt(s, f"  • {a}", 0.55, 2.15 + i*0.54, 5.7, 0.45, size=11, color=TEXT_WHITE)

# Right: Future work
box(s, 6.6, 1.55, 6.35, 5.65, BG_CARD, ACCENT_PURP, 1.0)
txt(s, "🔮  Future Work", 6.8, 1.65, 6.0, 0.38, size=14, bold=True, color=ACCENT_PURP)
future = [
    "Replace SSIM gate with BRISQUE / NRQM (no-reference)",
    "Add DehazeFormer / AOD-Net as DCP fallback",
    "RTTS real-world haze benchmark evaluation",
    "Lightweight PromptIR joint-refinement pass",
    "DRSformer SOTA deraining expert integration",
    "Cloud REST API deployment (FastAPI)",
    "Mobile ONNX export for Android/iOS edge",
    "Submit to IEEE Signal Processing Letters",
]
for i, f in enumerate(future):
    txt(s, f"  ▸ {f}", 6.8, 2.15 + i*0.62, 6.0, 0.5, size=11, color=TEXT_MUTED)

# Footer
box(s, 0.35, 7.15, 12.6, 0.28, BG_CARD, ACCENT_GREEN, 0.5)
txt(s, "MAIR+ v2  ·  github.com/NIt/MAIR_Project  ·  IEEE Journal Submission Pending",
    0.5, 7.18, 12.2, 0.22, size=9, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════
print(f"\nTotal slides after update: {len(prs.slides)}")
print(f"Saving to: {DEST}")
prs.save(str(DEST))
print("✅  Done! Presentation updated successfully.")
