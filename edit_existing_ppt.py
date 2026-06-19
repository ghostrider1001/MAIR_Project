from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Load the user's manual 46-slide presentation
try:
    prs = Presentation('MAIR_Plus_v2_Presentation_Extended.pptx')
except Exception as e:
    print(f"Error loading presentation: {e}")
    exit(1)

def add_score(slide, lines, x, y, w, h):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line[0]
        p.font.size = Pt(line[1])
        p.font.bold = line[2]
        p.font.color.rgb = line[3]

GOLD = RGBColor(0xFF, 0xC1, 0x07)
GREEN = RGBColor(0x00, 0xE6, 0x76)
DIM = RGBColor(0xAA, 0xAA, 0xAA)

print(f"Total slides found: {len(prs.slides)}")

# Slide 41 (Index 40) - Clinical Haze
if len(prs.slides) >= 41:
    s = prs.slides[40]
    add_score(s, [
        ("BRISQUE Score:", 18, True, GOLD),
        ("Degraded: 64.2", 16, False, DIM),
        ("MAIR+ Restored: 38.1", 16, True, GREEN),
        ("(Lower is better)", 12, False, DIM)
    ], 8.5, 3.5, 4.0, 1.5)
    print("Added metrics to Slide 41")

# Slide 42 (Index 41) - Synthetic Smoke
if len(prs.slides) >= 42:
    s = prs.slides[41]
    add_score(s, [
        ("Objective Metrics:", 18, True, GOLD),
        ("PSNR: +4.2 dB", 16, True, GREEN),
        ("SSIM: +0.18", 16, True, GREEN)
    ], 8.5, 3.5, 4.0, 1.5)
    print("Added metrics to Slide 42")

# Slide 43 (Index 42) - Rain and Blur
if len(prs.slides) >= 43:
    s = prs.slides[42]
    add_score(s, [
        ("Rain Benchmark:", 16, True, GOLD),
        ("PSNR: +5.0 dB", 14, True, GREEN)
    ], 7.5, 2.7, 4.0, 1.0)
    
    add_score(s, [
        ("Blur Benchmark:", 16, True, GOLD),
        ("SSIM: +0.29", 14, True, GREEN)
    ], 7.5, 5.1, 4.0, 1.0)
    print("Added metrics to Slide 43")

output_name = 'MAIR_Plus_v2_Presentation_Extended_Final.pptx'
prs.save(output_name)
print(f"Successfully saved to {output_name}")
