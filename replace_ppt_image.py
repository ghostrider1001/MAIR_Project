from pptx import Presentation
from pptx.util import Inches

print("Loading presentation...")
prs = Presentation('MAIR_Plus_v2_Presentation_Extended_Final.pptx')

# Slide 42 is at index 41
if len(prs.slides) >= 42:
    slide = prs.slides[41]
    
    # Remove existing picture
    pics_to_delete = []
    for shape in slide.shapes:
        if hasattr(shape, "image"):  # Identifies picture shapes
            pics_to_delete.append(shape)
            
    for p in pics_to_delete:
        sp = p._element
        sp.getparent().remove(sp)

    # Add the newly generated better image at the exact same coordinates
    img_path = r'outputs\comparison_grids\medical_synthetic_smoke_comparison.png'
    slide.shapes.add_picture(img_path, Inches(0.6), Inches(2.3), height=Inches(4.5))
    print("Image replaced on Slide 42!")

output_name = 'MAIR_Plus_v2_Presentation_Extended_Final_v2.pptx'
prs.save(output_name)
print(f"Presentation saved successfully to {output_name}")
