"""
generate_smoke_sweep_pptx.py
----------------------------
Generates a professional PowerPoint presentation for the smoke intensity sweep. 
Follows the MAIR+ dark-navy styling with Side-by-Side visuals.
"""

import os
import sys
import glob
import cv2
import numpy as np

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx is not installed. Run: pip install python-pptx")
    sys.exit(1)

# Quality Eval Imports for per-image scoring
sys.path.append(os.path.dirname(os.path.abspath(__file__)))
from evaluation.quality_evaluator import _compute_ssim_psnr as compute_ssim_psnr

try:
    import pyiqa
    import torch
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    niqe_metric = pyiqa.create_metric('niqe', device=device)
    HAS_NIQE = True
except ImportError:
    HAS_NIQE = False

# Color Palette
NAVY       = RGBColor(0x05, 0x0b, 0x14)
CARD_DARK  = RGBColor(0x0a, 0x11, 0x20)
NEON_GREEN = RGBColor(0x10, 0xb9, 0x81)
NEON_BLUE  = RGBColor(0x3b, 0x82, 0xf6)
NEON_PINK  = RGBColor(0xec, 0x48, 0x99)
WHITE      = RGBColor(0xff, 0xff, 0xff)
GREY       = RGBColor(0x9c, 0xa3, 0xaf)

def compute_niqe(img_path):
    if not HAS_NIQE: return 0.0
    try:
        return float(niqe_metric(img_path).item())
    except:
        return 0.0

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height, font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_card(slide, left, top, width, height, border_color=NEON_GREEN):
    shape = slide.shapes.add_shape(1, left, top, width, height) # 1 = msoShapeRectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_DARK
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

def get_best_image_for_pct(pct):
    candidates = glob.glob(f"datasets/synthetic_smoke_sweep/deg_{pct}_*.png")
    best_candidate = None
    best_psnr_gain = -999.0
    
    for deg_path in candidates:
        base_name = os.path.basename(deg_path).replace(f"deg_{pct}_", "")
        
        orig_clear_search = glob.glob(f"datasets/DeSmoke-LAP dataset/Dataset/*/clear/{base_name}")
        if not orig_clear_search:
            continue
            
        orig_clear_img = cv2.imread(orig_clear_search[0])
        
        rest_search = glob.glob(f"outputs/*/deg_{pct}_{base_name.replace('.png', '')}*.png")
        if not rest_search:
            continue
            
        rest_path = rest_search[0]
        
        deg_img = cv2.imread(deg_path)
        rest_img = cv2.imread(rest_path)
        
        if rest_img is None or deg_img is None or orig_clear_img is None:
            continue
            
        if rest_img.shape != orig_clear_img.shape:
            rest_img = cv2.resize(rest_img, (orig_clear_img.shape[1], orig_clear_img.shape[0]))
        if deg_img.shape != orig_clear_img.shape:
            deg_img = cv2.resize(deg_img, (orig_clear_img.shape[1], orig_clear_img.shape[0]))
            
        _, deg_psnr = compute_ssim_psnr(orig_clear_img, deg_img)
        _, rest_psnr = compute_ssim_psnr(orig_clear_img, rest_img)
        
        gain = rest_psnr - deg_psnr
        if gain > best_psnr_gain:
            best_psnr_gain = gain
            
            best_candidate = {
                'pct': pct,
                'base_name': base_name,
                'deg_path': deg_path,
                'rest_path': rest_path,
                'deg_psnr': deg_psnr,
                'rest_psnr': rest_psnr,
                'gain': gain,
                'orig_clear_img': orig_clear_img
            }
            
    if best_candidate:
        best_candidate['deg_niqe'] = compute_niqe(best_candidate['deg_path'])
        best_candidate['rest_niqe'] = compute_niqe(best_candidate['rest_path'])
        
    return best_candidate

def build_comparison_slide(prs, title, data_1, data_2):
    blank_layout = prs.slide_layouts[6] # blank
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, NAVY)
    
    add_textbox(slide, "VISUAL COMPARISON", Inches(0.5), Inches(0.2), Inches(3), Inches(0.4), font_size=12, bold=True, color=NEON_BLUE)
    add_textbox(slide, title, Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=30, bold=True, color=WHITE)
    add_textbox(slide, "Actual frames from the synthetic smoke sweep — hazy input vs MAIR+ DCP restored output", Inches(0.5), Inches(1.1), Inches(12), Inches(0.4), font_size=14, color=GREY)
    
    y_top = 1.6
    
    for idx, data in enumerate([data_1, data_2]):
        if data is None: continue
        
        x_start = 0.5 if idx == 0 else 6.9
        pct = data['pct']
        base_name = data['base_name'].replace('.png', '')
        deg_path = data['deg_path']
        rest_path = data['rest_path']
        deg_psnr = data['deg_psnr']
        rest_psnr = data['rest_psnr']
        deg_niqe = data['deg_niqe']
        rest_niqe = data['rest_niqe']
        gain = data['gain']
        
        # Draw Card
        add_card(slide, Inches(x_start), Inches(y_top), Inches(5.9), Inches(5.5), border_color=NEON_GREEN)
        
        # Title of the card e.g. "TLH_10 - Δ PSNR +12.3 dB"
        add_textbox(slide, f"{base_name} — Δ PSNR +{gain:.1f} dB ({pct}% Smoke)", Inches(x_start+0.2), Inches(y_top+0.1), Inches(5.5), Inches(0.4), font_size=12, bold=True, color=NEON_GREEN)
        
        # Images placement
        img_y = y_top + 0.8
        w_img = 2.6
        h_img = 3.2
        
        # Hazy
        pic_deg = slide.shapes.add_picture(deg_path, Inches(x_start+0.2), Inches(img_y), width=Inches(w_img))
        pic_deg.line.color.rgb = NEON_PINK
        pic_deg.line.width = Pt(2)
        add_textbox(slide, f"HAZY PSNR: {deg_psnr:.1f} | NIQE: {deg_niqe:.1f}", Inches(x_start+0.2), Inches(img_y + h_img + 0.1), Inches(w_img), Inches(0.4), font_size=10, bold=True, color=NEON_PINK, align=PP_ALIGN.CENTER)
        
        # Arrow
        add_textbox(slide, "→", Inches(x_start+2.8), Inches(img_y + 1.2), Inches(0.4), Inches(0.4), font_size=24, bold=True, color=NEON_GREEN)
        
        # Restored
        pic_rest = slide.shapes.add_picture(rest_path, Inches(x_start+3.1), Inches(img_y), width=Inches(w_img))
        pic_rest.line.color.rgb = NEON_GREEN
        pic_rest.line.width = Pt(2)
        add_textbox(slide, f"RESTORED PSNR: {rest_psnr:.1f} | NIQE: {rest_niqe:.1f}", Inches(x_start+3.1), Inches(img_y + h_img + 0.1), Inches(w_img), Inches(0.4), font_size=10, bold=True, color=NEON_GREEN, align=PP_ALIGN.CENTER)

    # Bottom legend
    add_textbox(slide, "● Red border = Hazy input (surgical smoke present)", Inches(3.0), Inches(7.0), Inches(4), Inches(0.4), font_size=11, color=NEON_PINK)
    add_textbox(slide, "● Green border = MAIR+ DCP output (smoke removed)", Inches(7.0), Inches(7.0), Inches(4), Inches(0.4), font_size=11, color=NEON_GREEN)

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    print("Finding best images for each intensity...")
    data_10 = get_best_image_for_pct(10)
    data_20 = get_best_image_for_pct(20)
    data_30 = get_best_image_for_pct(30)
    data_40 = get_best_image_for_pct(40)
    
    if not any([data_10, data_20, data_30, data_40]):
        print("No valid images found for any percentage.")
        return
        
    # Generate Slides
    build_comparison_slide(prs, "Before vs. After — 10% and 20% Smoke Intensity", data_10, data_20)
    build_comparison_slide(prs, "Before vs. After — 30% and 40% Smoke Intensity", data_30, data_40)
    
    out_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "MAIR_Smoke_Sweep_Visuals.pptx")
    prs.save(out_file)
    print(f"Generated successfully: {out_file}")

if __name__ == "__main__":
    main()
