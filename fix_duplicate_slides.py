"""
Remove the last N slides from MAIR_Plus_v2_Presentation_Final.pptx
Usage: python fix_duplicate_slides.py
"""
from pptx import Presentation
from pptx.oxml.ns import qn
from pathlib import Path

SLIDES_TO_REMOVE = 15  # remove the duplicate last 15 slides

path = Path("MAIR_Plus_v2_Presentation_Final.pptx")
prs = Presentation(str(path))

total_before = len(prs.slides)
print(f"Slides before fix: {total_before}")

if total_before < SLIDES_TO_REMOVE:
    print("ERROR: Not enough slides to remove. Aborting.")
    exit(1)

xml_slides = prs.slides._sldIdLst

# Collect rIds of slides to remove (last N), then remove backwards
# to avoid index shifting problems
for i in range(total_before - 1, total_before - SLIDES_TO_REMOVE - 1, -1):
    rId = xml_slides[i].get(qn('r:id'))
    # Drop the relationship from the presentation part
    prs.part.drop_rel(rId)
    # Remove the slide reference from the XML list
    del xml_slides[i]

total_after = len(prs.slides)
print(f"Slides after fix:  {total_after}")
print(f"Removed:           {total_before - total_after} slides")

prs.save(str(path))
print(f"Saved: {path}")
print("Done!")
